from pptx import Presentation
from pptx.util import Inches

# Создание презентации
prs = Presentation()

# Слайд 1: Титульный слайд
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Разработка мобильной платформы для обучения"
subtitle.text = "Студент: Волков А.О.\nГруппа: СМ7-74Б\nНаучный руководитель: Калиниченко С.В.\nМГТУ им. Н.Э. Баумана"

# Слайд 2: Актуальность работы
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Актуальность задачи"
content.text = (
    "- Необходимость улучшения подготовки специалистов в области робототехники.\n"
    "- Ограниченность доступных и эффективных образовательных платформ.\n"
    "- Высокая стоимость и сложность освоения существующих решений."
)

# Слайд 3: Цели и задачи
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Цели и задачи"
content.text = (
    "Цель:\n"
    "Разработка мобильной платформы для образовательных целей с возможностью изучения сложных алгоритмов управления и навигации.\n\n"
    "Задачи:\n"
    "1. Анализ существующих мобильных платформ.\n"
    "2. Разработка конструкции и системы управления.\n"
    "3. Подбор комплектующих элементов.\n"
    "4. Разработка алгоритмов SLAM-навигации."
)

# Слайд 4: Обзор существующих решений
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Обзор существующих решений"
content.text = (
    "Сравнение мобильных платформ:\n"
    "- KUKA youBot: Высокая маневренность, но дорогостоящая.\n"
    "- DJI RoboMaster EP: Поддержка Python и Scratch, но отсутствие лидара.\n"
    "- MYAGV: Компактная и доступная, с поддержкой ROS."
)

# Слайд 5: Технические требования
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Технические требования"
content.text = (
    "Основные характеристики:\n"
    "- Автономность: 2 часа.\n"
    "- Максимальная скорость: 1 м/с.\n"
    "- Наличие лидара.\n"
    "- Возможность подключения дополнительной периферии.\n"
    "- Грузоподъемность: не менее 8 кг."
)

# Слайд 6: Структура системы управления
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Структура системы управления"
content.text = (
    "Трехуровневая система управления:\n"
    "1. Верхний уровень: Управляющий компьютер (Raspberry Pi 5).\n"
    "2. Средний уровень: Микроконтроллер (STM32F407VET6).\n"
    "3. Нижний уровень: Приводы, лидар, датчики."
)

# Слайд 7: Подбор компонентов
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.shapes.placeholders[1].text_frame
title.text = "Подбор компонентов"
content.text = (
    "Основные компоненты:\n"
    "- Электродвигатель: CHP-42GP-775 ABHL.\n"
    "- Энкодер: Магнитный энкодер с разрешением 425 импульсов/оборот.\n"
    "- Лидар: RP-Lidar A1.\n"
    "- Камера: Raspberry Pi Camera Module 2."
)

# Сохранение презентации
prs.save('mobile_platform_presentation.pptx')